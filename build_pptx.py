# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (aligne sur le theme Power BI FEI_Humanitarian) ----
NAVY = RGBColor(0x14, 0x1B, 0x2E)       # fond slides sombres
INK = RGBColor(0x25, 0x23, 0x1F)        # texte sur fond clair
CREAM = RGBColor(0xFD, 0xFC, 0xFA)      # fond slides claires (blanc casse tres leger, coherent theme PBI)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x8A, 0x84, 0x78)
PRIMARY = RGBColor(0x00, 0x72, 0xB2)    # bleu institutionnel
TEAL = RGBColor(0x00, 0x9E, 0x73)
AMBER = RGBColor(0xE8, 0xA3, 0x3D)
RED = RGBColor(0xD6, 0x48, 0x3F)
GREEN = RGBColor(0x1A, 0xAB, 0x40)
CARD_BG = RGBColor(0xF4, 0xF1, 0xEA)

FONT_BODY = "Calibri"
FONT_HEAD = "Cambria"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    bgfill = s.background.fill
    bgfill.solid()
    bgfill.fore_color.rgb = bg
    return s


def textbox(slide, l, t, w, h, text, size=14, color=INK, bold=False, italic=False,
            font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return tb


def multiline(slide, l, t, w, h, lines, size=13, color=INK, font=FONT_BODY,
              align=PP_ALIGN.LEFT, space_after=8, bullet=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = 1.12
        if bullet:
            pPr = p._pPr
            if pPr is None:
                pPr = p._p.get_or_add_pPr()
            buChar = pPr.makeelement(qn('a:buChar'), {'char': '–'})
            buFont = pPr.makeelement(qn('a:buFont'), {'typeface': font})
            pPr.append(buFont)
            pPr.append(buChar)
            p_indent = Emu(Inches(0.22))
            pPr.set('marL', str(p_indent))
            pPr.set('indent', str(-p_indent))
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
    return tb


def rounded_card(slide, l, t, w, h, fill=WHITE, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.adjustments[0] = 0.06
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shadow = shp.shadow
    shadow.inherit = False
    return shp


def stat_card(slide, l, t, w, h, value, label, accent=PRIMARY):
    card = rounded_card(slide, l, t, w, h, fill=WHITE)
    value_size = 26 if len(value) <= 9 else 19
    textbox(slide, l + 0.18, t + 0.12, w - 0.36, h * 0.55, value, size=value_size, color=accent,
            bold=True, font=FONT_HEAD, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.BOTTOM)
    textbox(slide, l + 0.18, t + h * 0.60, w - 0.36, h * 0.35, label, size=10.5, color=GREY,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.05)
    return card


def slide_header(slide, kicker, title, dark=False):
    kcolor = TEAL if dark else PRIMARY
    tcolor = WHITE if dark else INK
    textbox(slide, 0.6, 0.35, 10, 0.4, kicker.upper(), size=12, color=kcolor, bold=True, font=FONT_BODY)
    textbox(slide, 0.6, 0.68, 11.5, 0.75, title, size=28, color=tcolor, bold=True, font=FONT_HEAD)


def page_number(slide, n, dark=False):
    textbox(slide, 12.6, 7.08, 0.6, 0.35, str(n), size=10, color=GREY if not dark else RGBColor(0x9A, 0xA3, 0xB5),
            align=PP_ALIGN.RIGHT)


def footer_source(slide, dark=False):
    textbox(slide, 0.6, 7.08, 8, 0.35,
            "Source : Base de données MEAL FEI 2021-2025 — Tableau de bord Power BI",
            size=9, color=GREY if not dark else RGBColor(0x9A, 0xA3, 0xB5), italic=True)


def add_chart(slide, l, t, w, h, chart_type, categories, series_name, values,
              color=PRIMARY, number_format='#,##0', title=None, horizontal=False):
    data = CategoryChartData()
    data.categories = categories
    data.add_series(series_name, values)
    gframe = slide.shapes.add_chart(chart_type, Inches(l), Inches(t), Inches(w), Inches(h), data)
    chart = gframe.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = number_format
    dl.number_format_is_linked = False
    dl.font.size = Pt(10)
    dl.font.bold = True
    dl.font.color.rgb = INK
    try:
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
    except Exception:
        pass
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    series.format.line.fill.background()
    cat_ax = chart.category_axis
    val_ax = chart.value_axis
    cat_ax.tick_labels.font.size = Pt(11)
    cat_ax.tick_labels.font.color.rgb = INK
    cat_ax.format.line.color.rgb = RGBColor(0xE4, 0xDF, 0xD5)
    val_ax.visible = False
    val_ax.has_major_gridlines = False
    cat_ax.has_major_gridlines = False
    if chart.chart_type in (XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS):
        series.format.line.color.rgb = color
        series.format.line.width = Pt(2.5)
        series.smooth = False
        series.marker.style = 8  # circle
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = color
    return chart


# ============================================================
# SLIDE 1 -- TITRE
# ============================================================
s = add_slide(bg=NAVY)
textbox(s, 0.9, 2.35, 11.5, 0.5, "RAPPORT D'ANALYSE — DIRECTION", size=15, color=TEAL, bold=True, font=FONT_BODY)
textbox(s, 0.9, 2.85, 11.5, 1.8,
        "Performance Programmatique, Financière,\nOpérationnelle et de Redevabilité",
        size=38, color=WHITE, bold=True, font=FONT_HEAD, line_spacing=1.05)
textbox(s, 0.9, 4.35, 11.5, 0.5, "Faso Espoir International (FEI) — Burkina Faso", size=18, color=RGBColor(0xCA, 0xDC, 0xFC), font=FONT_BODY)
textbox(s, 0.9, 5.05, 11.5, 0.4, "Période analysée : 2021–2025   |   Devise de référence : USD   |   18 juillet 2026",
        size=13, color=RGBColor(0x9A, 0xA3, 0xB5), font=FONT_BODY)
card = rounded_card(s, 0.9, 6.35, 8.6, 0.6, fill=RGBColor(0x24, 0x2C, 0x44))
textbox(s, 1.15, 6.35, 8.2, 0.6,
        "CONFIDENTIEL — données bénéficiaires anonymisées et agrégées",
        size=11, color=RGBColor(0xE8, 0xA3, 0x3D), bold=True, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# SLIDE 2 -- CONTEXTE & METHODOLOGIE
# ============================================================
s = add_slide()
slide_header(s, "Contexte", "Une ONG humanitaire active depuis 2021 au Burkina Faso")
multiline(s, 0.6, 1.75, 6.0, 4.6, [
    "ONG internationale (ONGI) intervenant dans le domaine MEAL — Suivi, "
    "Évaluation, Redevabilité et Apprentissage",
    "12 bases terrain et 11 secteurs d'intervention humanitaire (sécurité "
    "alimentaire, nutrition, EHA, protection…)",
    "Coordination avec le cluster OCHA Burkina Faso",
    "Objet du rapport : éclairer la Direction sur 4 axes — programme, "
    "finance, opérations, redevabilité",
], size=15, color=INK)

card = rounded_card(s, 6.95, 1.75, 5.8, 4.6, fill=CARD_BG)
textbox(s, 7.2, 1.95, 5.3, 0.4, "MÉTHODOLOGIE", size=13, color=PRIMARY, bold=True, font=FONT_BODY)
multiline(s, 7.2, 2.45, 5.3, 3.8, [
    "Base de données MEAL consolidée : 25 tables (bénéficiaires, "
    "activités, décaissements, cadre logique, plaintes, stocks, partenaires)",
    "Collecte terrain via KoboToolbox / ODK",
    "Période : 2021–2025 (5 années fiscales complètes)",
    "Nettoyage, modélisation et calcul de 31 indicateurs clés",
    "Données nominatives des bénéficiaires exclues de l'analyse",
], size=13.5, color=INK)
page_number(s, 2)
footer_source(s)

# ============================================================
# SLIDE 3 -- PERFORMANCE PROGRAMMATIQUE
# ============================================================
s = add_slide()
slide_header(s, "Performance programmatique", "Une couverture ciblée, un cadre logique inégal")

stat_card(s, 0.6, 1.75, 2.55, 1.4, "1 372 / 2 000", "Bénéficiaires actifs (68,6 %)", accent=PRIMARY)
stat_card(s, 3.3, 1.75, 2.55, 1.4, "70,6 %", "Réalisation cadre logique (volume)", accent=TEAL)
stat_card(s, 0.6, 3.3, 2.55, 1.4, "20,6 %", "Objectifs individuels « Atteint »", accent=AMBER)
stat_card(s, 3.3, 3.3, 2.55, 1.4, "51 / 100", "Score de vulnérabilité moyen", accent=RED)

add_chart(s, 6.2, 1.75, 6.5, 4.35, XL_CHART_TYPE.BAR_CLUSTERED,
          ["Boucle du Mouhoun", "Centre", "Nord", "Centre-Nord", "Sahel", "Est"],
          "Bénéficiaires", [151, 194, 220, 267, 267, 305], color=PRIMARY,
          number_format='#,##0')
textbox(s, 6.2, 1.42, 6.5, 0.3, "Bénéficiaires enregistrés par région (top 6)", size=12, color=GREY, italic=True)

page_number(s, 3)
footer_source(s)

# ============================================================
# SLIDE 4 -- SUIVI FINANCIER
# ============================================================
s = add_slide()
slide_header(s, "Suivi financier", "Une exécution budgétaire solide et stable")

stat_card(s, 0.6, 1.75, 2.55, 1.4, "91,7 %", "Taux d'exécution budgétaire", accent=GREEN)
stat_card(s, 3.3, 1.75, 2.55, 1.4, "24,11 M $", "Total décaissé (2021-2025)", accent=PRIMARY)
stat_card(s, 0.6, 3.3, 2.55, 1.4, "~20 %", "Frais de gestion (stable 5 ans)", accent=TEAL)
stat_card(s, 3.3, 3.3, 2.55, 1.4, "10,3 %", "Décaissements en attente / rejetés", accent=RED)

add_chart(s, 6.2, 1.75, 6.5, 4.35, XL_CHART_TYPE.LINE_MARKERS,
          ["2021", "2022", "2023", "2024", "2025"],
          "Valeur de l'assistance (USD)",
          [4341530, 4238088, 4191723, 4390093, 4470656], color=PRIMARY,
          number_format='#,##0')
textbox(s, 6.2, 1.42, 6.5, 0.3, "Valeur de l'assistance délivrée par année (USD)", size=12, color=GREY, italic=True)

page_number(s, 4)
footer_source(s)

# ============================================================
# SLIDE 5 -- OPERATIONS & APPROVISIONNEMENT
# ============================================================
s = add_slide()
slide_header(s, "Opérations & approvisionnement", "Une logistique globalement maîtrisée, un risque de rupture à traiter")

stat_card(s, 0.6, 1.75, 2.55, 1.4, "82,2 %", "Activités réalisées (/15 000)", accent=PRIMARY)
stat_card(s, 3.3, 1.75, 2.55, 1.4, "3,3 j.", "Délai moyen de livraison", accent=TEAL)
stat_card(s, 0.6, 3.3, 2.55, 1.4, "12,1 %", "Livraisons en retard", accent=AMBER)
stat_card(s, 3.3, 3.3, 2.55, 1.4, "11 / 12", "Bases en rupture de stock imminente", accent=RED)

add_chart(s, 6.2, 1.75, 6.5, 4.35, XL_CHART_TYPE.COLUMN_CLUSTERED,
          ["Réalisée", "Planifiée", "Reportée", "Annulée"],
          "Activités", [12325, 1218, 853, 604], color=TEAL,
          number_format='#,##0')
textbox(s, 6.2, 1.42, 6.5, 0.3, "Répartition des 15 000 activités par statut", size=12, color=GREY, italic=True)

page_number(s, 5)
footer_source(s)

# ============================================================
# SLIDE 6 -- REDEVABILITE & QUALITE
# ============================================================
s = add_slide()
slide_header(s, "Redevabilité & qualité", "Un mécanisme actif, un signal fort sur l'équité du ciblage")

stat_card(s, 0.6, 1.75, 2.55, 1.4, "56,7 %", "Plaintes résolues (/2 000)", accent=GREEN)
stat_card(s, 3.3, 1.75, 2.55, 1.4, "0,80 / 1", "Score performance partenaires", accent=PRIMARY)
stat_card(s, 0.6, 3.3, 2.55, 1.4, "21,8 %", "Plaintes liées au ciblage/éligibilité", accent=RED)
stat_card(s, 3.3, 3.3, 2.55, 1.4, "18,1 j.", "Délai moyen de rapportage partenaires", accent=AMBER)

add_chart(s, 6.2, 1.75, 6.5, 4.35, XL_CHART_TYPE.BAR_CLUSTERED,
          ["Qualité assistance", "Perte docs/statut", "Retard distribution", "Discrimination ciblage", "Éligibilité contestée"],
          "Plaintes", [202, 203, 209, 210, 218], color=RED,
          number_format='#,##0')
textbox(s, 6.2, 1.42, 6.5, 0.3, "Principaux motifs de plainte (top 5)", size=12, color=GREY, italic=True)

page_number(s, 6)
footer_source(s)

# ============================================================
# SLIDE 7 -- RECOMMANDATIONS (finale)
# ============================================================
s = add_slide(bg=NAVY)
textbox(s, 0.6, 0.5, 11.5, 0.4, "RECOMMANDATIONS", size=15, color=TEAL, bold=True, font=FONT_BODY)
textbox(s, 0.6, 0.85, 11.8, 0.75, "Quatre priorités pour la Direction", size=30, color=WHITE, bold=True, font=FONT_HEAD)

recos = [
    ("1", RED, "Auditer l'équité du ciblage", "32 % des plaintes portent sur l'éligibilité ou une discrimination perçue dans le ciblage."),
    ("2", RED, "Sécuriser la chaîne d'approvisionnement", "11 bases sur 12 ont connu une rupture de stock imminente — revoir les seuils et délais."),
    ("3", RED, "Piloter le cadre logique objectif par objectif", "43 % des lignes d'objectifs sont « Non Atteint » malgré 71 % de réalisation globale."),
    ("4", AMBER, "Fiabiliser le circuit de décaissement", "10,3 % des décaissements restent en attente ou rejetés — analyser les causes."),
]

x = 0.6
w = 2.95
for i, (num, color, title, text) in enumerate(recos):
    card = rounded_card(s, x, 1.95, w, 4.35, fill=RGBColor(0x1D, 0x25, 0x3D))
    textbox(s, x + 0.25, 2.15, w - 0.5, 0.9, num, size=40, color=color, bold=True, font=FONT_HEAD)
    textbox(s, x + 0.25, 3.05, w - 0.5, 1.1, title, size=15.5, color=WHITE, bold=True, font=FONT_BODY, line_spacing=1.05)
    textbox(s, x + 0.25, 4.15, w - 0.5, 2.0, text, size=11.5, color=RGBColor(0xC7, 0xCE, 0xDD), font=FONT_BODY, line_spacing=1.15)
    x += w + 0.25

textbox(s, 0.6, 6.75, 11.8, 0.5,
        "À maintenir : discipline budgétaire (91,7 % d'exécution) et mécanisme de plaintes actif (56,7 % de résolution).",
        size=12, color=GREEN, italic=True, bold=True)

prs.save("Presentation_Direction_FEI_2021-2025.pptx")
print("OK")
